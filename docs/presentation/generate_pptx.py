"""
Crypto Market Monitor - Presentation PPTX v2
EFREI branding officiel : navy #163767, rose #ff43b8, bleu #3653a0
Principes presentation-design : massive headlines, section labels, progress bar, minimal text
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, math

BASE = os.path.dirname(os.path.abspath(__file__))
OUT  = os.path.join(BASE, "crypto-market-monitor.pptx")
LOGO = os.path.join(BASE, "efrei-logo.png")

# ── EFREI Brand palette ──────────────────────────────────────────
BG         = RGBColor(0x0B, 0x1B, 0x34)   # near-black navy
NAVY       = RGBColor(0x16, 0x37, 0x67)   # --bleuFonce
ROSE       = RGBColor(0xFF, 0x43, 0xB8)   # --rose  (accent principal EFREI)
BLUE       = RGBColor(0x36, 0x53, 0xA0)   # --bleuClair
BLUE2      = RGBColor(0x33, 0x7A, 0xB6)   # --bleuTresClair
PURPLE     = RGBColor(0x95, 0x56, 0x9E)   # --roseFonce
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
MUTED      = RGBColor(0x94, 0xA3, 0xB8)
CARD       = RGBColor(0x16, 0x37, 0x67)   # navy card bg
CARD2      = RGBColor(0x1E, 0x2E, 0x4A)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
DARK_CARD  = RGBColor(0x0D, 0x1F, 0x3C)

# Section accent colors (presentation-design spec)
SEC = {
    "intro":  ROSE,
    "archi":  BLUE,
    "kafka":  PURPLE,
    "svc":    BLUE2,
    "dash":   ROSE,
    "secu":   RGBColor(0xF8, 0x71, 0x71),
    "i18n":   RGBColor(0x34, 0xD3, 0x99),
    "deploy": BLUE,
    "bilan":  ROSE,
}

TOTAL_SLIDES = 12

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return s


def rect(slide, l, t, w, h, color=CARD, alpha=None):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def txt(slide, text, l, t, w, h,
        size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.italic = italic
    return tb


def progress_bar(slide, num, color=ROSE):
    frac = num / TOTAL_SLIDES
    w = 13.33 * frac
    bar = slide.shapes.add_shape(1, Inches(0), Inches(7.44), Inches(13.33), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0x16, 0x37, 0x67)
    bar.line.fill.background()
    bar2 = slide.shapes.add_shape(1, Inches(0), Inches(7.44), Inches(w), Inches(0.06))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = color
    bar2.line.fill.background()


def section_label(slide, label, color=ROSE):
    txt(slide, label.upper(), 0.45, 0.2, 8, 0.35,
        size=9, bold=True, color=color, align=PP_ALIGN.LEFT)


def slide_num(slide, num, color=MUTED):
    txt(slide, f"{num} / {TOTAL_SLIDES}", 12.5, 0.18, 0.8, 0.3,
        size=9, color=color, align=PP_ALIGN.RIGHT)


def headline(slide, text, l=0.5, t=0.7, w=12.3, h=2.0, size=44):
    txt(slide, text, l, t, w, h, size=size, bold=True, color=WHITE)


def add_logo(slide):
    if os.path.exists(LOGO):
        try:
            slide.shapes.add_picture(LOGO, Inches(11.8), Inches(0.15), width=Inches(1.25))
        except Exception:
            pass


def bullet(slide, items, l, t, w, size=13, color=WHITE, spacing=0.44):
    for i, item in enumerate(items):
        txt(slide, f"  {item}", l, t + i * spacing, w, spacing + 0.06,
            size=size, color=color)


# ═══════════════════════════════════════════════════════════════════
# Slide 1 — TITRE  (Big Statement)
# ═══════════════════════════════════════════════════════════════════
s1 = new_slide()
# Gradient band left (EFREI rose)
band = s1.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
band.fill.solid(); band.fill.fore_color.rgb = ROSE; band.line.fill.background()
# Top line accent
top = s1.shapes.add_shape(1, Inches(0.4), Inches(0), Inches(12.93), Inches(0.06))
top.fill.solid(); top.fill.fore_color.rgb = ROSE; top.line.fill.background()

txt(s1, "CRYPTO MONITOR", 0.8, 0.8, 11.5, 1.5, size=56, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txt(s1, "Surveillance temps reel des marches crypto", 0.8, 2.45, 10, 0.7,
    size=22, color=MUTED, align=PP_ALIGN.LEFT)

# Pink divider line
div = s1.shapes.add_shape(1, Inches(0.8), Inches(3.3), Inches(3), Inches(0.05))
div.fill.solid(); div.fill.fore_color.rgb = ROSE; div.line.fill.background()

txt(s1, "Apache Kafka  -  Node.js  -  TypeScript  -  Socket.IO", 0.8, 3.5, 10, 0.5,
    size=14, color=RGBColor(0xFF, 0x43, 0xB8), italic=True)

# Author card
rect(s1, 0.8, 4.5, 7.0, 1.8, DARK_CARD)
txt(s1, "Adam Beloucif  &  Emilien Morice", 1.0, 4.65, 6.6, 0.55,
    size=16, bold=True, color=WHITE)
txt(s1, "M1 Data Engineering & IA  -  EFREI Paris", 1.0, 5.25, 6.6, 0.45,
    size=12, color=ROSE)
txt(s1, "Module Real-Time Engineering  -  2025-2026", 1.0, 5.68, 6.6, 0.38,
    size=11, color=MUTED)

add_logo(s1)
progress_bar(s1, 1)

# ═══════════════════════════════════════════════════════════════════
# Slide 2 — CONTEXTE  (Full statement + metrics boxes)
# ═══════════════════════════════════════════════════════════════════
s2 = new_slide()
section_label(s2, "contexte & problematique", SEC["intro"])
slide_num(s2, 2, SEC["intro"])
add_logo(s2)

headline(s2, "Pourquoi surveiller\nles marches en temps reel ?", t=0.5, size=36)

# 3 metric boxes
metrics = [("100M+", "transactions/jour\nBinance"), ("< 1s", "latence cible\ndetection"), ("3", "sources\nBinance + Coinbase")]
for i, (val, label) in enumerate(metrics):
    x = 0.5 + i * 4.3
    rect(s2, x, 3.55, 3.8, 1.8, CARD)
    div2 = s2.shapes.add_shape(1, Inches(x), Inches(3.55), Inches(3.8), Inches(0.06))
    div2.fill.solid(); div2.fill.fore_color.rgb = ROSE; div2.line.fill.background()
    txt(s2, val,   x+0.15, 3.65, 3.5, 0.75, size=30, bold=True, color=ROSE)
    txt(s2, label, x+0.15, 4.38, 3.5, 0.85, size=12, color=MUTED)

txt(s2, "Kafka decouples ingestion from processing - N consumers, zero data loss",
    0.5, 5.65, 12.3, 0.45, size=12, color=MUTED, italic=True)
progress_bar(s2, 2)

# ═══════════════════════════════════════════════════════════════════
# Slide 3 — ARCHITECTURE  (Split layout)
# ═══════════════════════════════════════════════════════════════════
s3 = new_slide()
section_label(s3, "architecture", SEC["archi"])
slide_num(s3, 3, SEC["archi"])
add_logo(s3)

headline(s3, "Pipeline\nKafka bout-en-bout", t=0.5, size=36)

pipeline = [
    ("WebSocket\nSources", PURPLE, "Binance + Coinbase"),
    ("Ingester\nNode.js", SEC["archi"], "normalisation"),
    ("Kafka\nKRaft 3.7", RGBColor(0x23,0x1F,0x20), "tampon"),
    ("Processor\nVWAP / SMA", SEC["archi"], "analytics"),
    ("API\nSocket.IO", ROSE, "live push"),
]
for i, (name, col, sub) in enumerate(pipeline):
    x = 0.35 + i * 2.55
    box = s3.shapes.add_shape(1, Inches(x), Inches(3.3), Inches(2.3), Inches(1.4))
    box.fill.solid(); box.fill.fore_color.rgb = col; box.line.fill.background()
    tf = box.text_frame; tf.text = name
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(11); run.font.bold = True
            run.font.color.rgb = WHITE
    txt(s3, sub, x, 4.82, 2.3, 0.3, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 4:
        arr = s3.shapes.add_shape(1, Inches(x+2.32), Inches(3.9), Inches(0.2), Inches(0.2))
        arr.fill.solid(); arr.fill.fore_color.rgb = MUTED; arr.line.fill.background()

rect(s3, 0.35, 5.35, 12.6, 1.75, DARK_CARD)
div3 = s3.shapes.add_shape(1, Inches(0.35), Inches(5.35), Inches(0.05), Inches(1.75))
div3.fill.solid(); div3.fill.fore_color.rgb = SEC["archi"]; div3.line.fill.background()
txt(s3, "Decoupling garanti", 0.6, 5.45, 12, 0.4, size=13, bold=True, color=WHITE)
txt(s3, "Les producteurs et consommateurs fonctionnent a des debits independants. Partition par symbole (BTC-USDT, ETH-USDT, BTC-USD). Retention 1h, relecture possible en cas de crash.",
    0.6, 5.88, 12.1, 1.1, size=11, color=MUTED)
progress_bar(s3, 3, SEC["archi"])

# ═══════════════════════════════════════════════════════════════════
# Slide 4 — KAFKA KRaft  (Full statement + key points)
# ═══════════════════════════════════════════════════════════════════
s4 = new_slide()
section_label(s4, "apache kafka", SEC["kafka"])
slide_num(s4, 4, SEC["kafka"])
add_logo(s4)

headline(s4, "KRaft : Kafka\nsans Zookeeper", t=0.5, size=40)

rect(s4, 0.5, 3.25, 5.7, 3.85, DARK_CARD)
div4 = s4.shapes.add_shape(1, Inches(0.5), Inches(3.25), Inches(0.07), Inches(3.85))
div4.fill.solid(); div4.fill.fore_color.rgb = SEC["kafka"]; div4.line.fill.background()
txt(s4, "Kafka 3.7 gere ses metadonnees via Raft\ninterne - 1 seul conteneur, 0 Zookeeper.",
    0.7, 3.35, 5.3, 1.0, size=13, color=WHITE)
txt(s4, "bitnami/kafka:3.7  (KRaft pre-configure)", 0.7, 4.4, 5.3, 0.4, size=11,
    color=SEC["kafka"], italic=True)
bullet(s4, [
    "Producteurs ecrivent a la vitesse du marche",
    "Consommateurs lisent a leur propre rythme",
    "Partition par symbole, ordre garanti",
    "Retention 1h - relecture en cas de crash",
], 0.7, 4.9, 5.2, size=12, spacing=0.46)

rect(s4, 6.8, 3.25, 6.1, 3.85, DARK_CARD)
div4b = s4.shapes.add_shape(1, Inches(6.8), Inches(3.25), Inches(0.07), Inches(3.85))
div4b.fill.solid(); div4b.fill.fore_color.rgb = BLUE; div4b.line.fill.background()
txt(s4, "Topics", 7.0, 3.35, 5.7, 0.4, size=14, bold=True, color=WHITE)
topics = [("crypto-trades", "Transactions brutes normalisees", "3 partitions, 1h"),
          ("crypto-metrics","Metriques calculees (VWAP...)", "3 partitions, 1h")]
for i, (topic, desc, conf) in enumerate(topics):
    y = 3.9 + i * 1.4
    rect(s4, 7.0, y, 5.7, 1.2, CARD)
    txt(s4, topic, 7.15, y+0.1, 5.3, 0.4, size=12, bold=True, color=SEC["kafka"])
    txt(s4, desc,  7.15, y+0.52, 5.3, 0.38, size=11, color=WHITE)
    txt(s4, conf,  7.15, y+0.86, 5.3, 0.28, size=9, color=MUTED, italic=True)
progress_bar(s4, 4, SEC["kafka"])

# ═══════════════════════════════════════════════════════════════════
# Slide 5 — INGESTER  (Split)
# ═══════════════════════════════════════════════════════════════════
s5 = new_slide()
section_label(s5, "service ingester", SEC["svc"])
slide_num(s5, 5, SEC["svc"])
add_logo(s5)

headline(s5, "Collecte\ndes flux WebSocket", t=0.5, size=36)

rect(s5, 0.4, 3.2, 5.9, 3.9, DARK_CARD)
div5 = s5.shapes.add_shape(1, Inches(0.4), Inches(3.2), Inches(0.07), Inches(3.9))
div5.fill.solid(); div5.fill.fore_color.rgb = SEC["svc"]; div5.line.fill.background()
txt(s5, "Sources", 0.65, 3.3, 5.5, 0.4, size=13, bold=True, color=WHITE)
bullet(s5, [
    "Binance : BTC/USDT + ETH/USDT (stream combine)",
    "Coinbase : BTC-USD (Advanced Trade WS)",
    "Interface Trade normalisee",
    "Cle Kafka = symbole",
], 0.65, 3.8, 5.5, size=12, spacing=0.46)
txt(s5, "{ exchange, symbol, price, quantity,\n  timestamp, tradeId, side, value }",
    0.7, 5.75, 5.5, 0.9, size=10, color=GREEN, italic=True)

rect(s5, 6.9, 3.2, 6.1, 3.9, DARK_CARD)
div5b = s5.shapes.add_shape(1, Inches(6.9), Inches(3.2), Inches(0.07), Inches(3.9))
div5b.fill.solid(); div5b.fill.fore_color.rgb = ROSE; div5b.line.fill.background()
txt(s5, "Resilience", 7.15, 3.3, 5.7, 0.4, size=13, bold=True, color=WHITE)
bullet(s5, [
    "Reconnexion auto - backoff exponentiel",
    "Max 10 tentatives, delai max 30s",
    "Heartbeat Coinbase gere",
    "Arret propre SIGTERM/SIGINT",
    "Vidage buffer Kafka avant arret",
], 7.15, 3.8, 5.7, size=12, spacing=0.46)
progress_bar(s5, 5, SEC["svc"])

# ═══════════════════════════════════════════════════════════════════
# Slide 6 — PROCESSOR  (Data/metrics layout)
# ═══════════════════════════════════════════════════════════════════
s6 = new_slide()
section_label(s6, "service processor", SEC["svc"])
slide_num(s6, 6, SEC["svc"])
add_logo(s6)

headline(s6, "Analytique\nen fenetre glissante", t=0.5, size=36)

rows = [
    ("VWAP",      "sum(prix x qte) / sum(qte)",  "1 min",    "Prix moyen pondere par volume"),
    ("SMA-20",    "moyenne 20 derniers prix",     "20 trades","Tendance lissee"),
    ("Z-Score",   "(val - mu) / sigma",           "1 min",    "Anomalie si > 2.5 sigma"),
    ("Var. 1min", "(prix - prix0) / prix0",       "1 min",    "Performance instantanee"),
    ("Var. 5min", "(prix - prix0) / prix0",       "5 min",    "Performance moyen terme"),
    ("High/Low",  "max / min des prix",           "1 min",    "Amplitude de la bougie"),
]
# Header row
rect(s6, 0.4, 3.2, 12.6, 0.42, CARD)
for cx, lbl in [(0.55,"Metrique"), (2.9,"Formule"), (7.0,"Fenetre"), (8.8,"Objectif")]:
    txt(s6, lbl, cx, 3.27, 3, 0.3, size=10, bold=True, color=ROSE)

for i, (name, formula, window, desc) in enumerate(rows):
    y = 3.65 + i * 0.55
    bg = DARK_CARD if i%2==0 else CARD
    rect(s6, 0.4, y, 12.6, 0.53, bg)
    txt(s6, name,    0.55, y+0.09, 2.2, 0.35, size=12, bold=True, color=WHITE)
    txt(s6, formula, 2.9,  y+0.09, 3.9, 0.35, size=10, color=GREEN)
    txt(s6, window,  7.0,  y+0.09, 1.6, 0.35, size=11, color=ROSE)
    txt(s6, desc,    8.8,  y+0.09, 4.0, 0.35, size=10, color=MUTED)

txt(s6, "Buffer circulaire : max 1000 trades / symbole, eviction TTL 10 minutes",
    0.55, 7.05, 12, 0.35, size=10, color=MUTED, italic=True)
progress_bar(s6, 6, SEC["svc"])

# ═══════════════════════════════════════════════════════════════════
# Slide 7 — API  (Code slide style)
# ═══════════════════════════════════════════════════════════════════
s7 = new_slide()
section_label(s7, "service api", SEC["svc"])
slide_num(s7, 7, SEC["svc"])
add_logo(s7)

headline(s7, "Express + Socket.IO\nREST & Push temps reel", t=0.5, size=34)

rect(s7, 0.4, 3.2, 5.9, 3.9, DARK_CARD)
div7 = s7.shapes.add_shape(1, Inches(0.4), Inches(3.2), Inches(0.07), Inches(3.9))
div7.fill.solid(); div7.fill.fore_color.rgb = SEC["svc"]; div7.line.fill.background()
txt(s7, "REST endpoints", 0.65, 3.3, 5.5, 0.4, size=13, bold=True, color=WHITE)
for i, ep in enumerate([
    "GET /api/health",
    "GET /api/metrics",
    "GET /api/metrics/:symbol",
    "GET /api/history/:symbol",
    "GET /api/anomalies",
]):
    txt(s7, ep, 0.7, 3.85 + i*0.55, 5.4, 0.45, size=12, color=GREEN)

rect(s7, 6.9, 3.2, 6.1, 3.9, DARK_CARD)
div7b = s7.shapes.add_shape(1, Inches(6.9), Inches(3.2), Inches(0.07), Inches(3.9))
div7b.fill.solid(); div7b.fill.fore_color.rgb = ROSE; div7b.line.fill.background()
txt(s7, "Evenements Socket.IO (push)", 7.15, 3.3, 5.7, 0.4, size=13, bold=True, color=WHITE)
for i, (ev, desc) in enumerate([
    ("metrics:update",   "Metriques calculees"),
    ("anomaly:detected", "Alerte z-score > seuil"),
    ("server:stats",     "Stats connexions / msg/s"),
    ("initial:state",    "Etat complet a la connexion"),
    ("subscribe:symbol", "Choix de la paire"),
]):
    y = 3.85 + i*0.55
    txt(s7, ev,   7.15, y, 3.0, 0.38, size=11, bold=True, color=ROSE)
    txt(s7, desc, 10.2, y, 2.6, 0.38, size=11, color=MUTED)
progress_bar(s7, 7, SEC["svc"])

# ═══════════════════════════════════════════════════════════════════
# Slide 8 — DASHBOARD  (Split)
# ═══════════════════════════════════════════════════════════════════
s8 = new_slide()
section_label(s8, "dashboard", SEC["dash"])
slide_num(s8, 8, SEC["dash"])
add_logo(s8)

headline(s8, "Interface\ntout-en-un live", t=0.5, size=40)

rect(s8, 0.4, 3.2, 5.9, 3.9, DARK_CARD)
div8 = s8.shapes.add_shape(1, Inches(0.4), Inches(3.2), Inches(0.07), Inches(3.9))
div8.fill.solid(); div8.fill.fore_color.rgb = ROSE; div8.line.fill.background()
txt(s8, "Composants", 0.65, 3.3, 5.5, 0.4, size=13, bold=True, color=WHITE)
bullet(s8, [
    "Prix hero - flash vert/rouge au changement",
    "Badges variation 1min et 5min",
    "4 stat-cards : Volume, Trades, High, Low",
    "Chart.js : ligne prix + SMA-20 superposee",
    "Barres volume par paire",
    "Flux anomalies + beep AudioContext",
], 0.65, 3.8, 5.5, size=12, spacing=0.44)

rect(s8, 6.9, 3.2, 6.1, 3.9, DARK_CARD)
div8b = s8.shapes.add_shape(1, Inches(6.9), Inches(3.2), Inches(0.07), Inches(3.9))
div8b.fill.solid(); div8b.fill.fore_color.rgb = BLUE; div8b.line.fill.background()
txt(s8, "Design & UX", 7.15, 3.3, 5.7, 0.4, size=13, bold=True, color=WHITE)
bullet(s8, [
    "Dark theme EFREI navy + rose #ff43b8",
    "Logos SVG : Bitcoin, Ethereum, Coinbase",
    "Barre gradient EFREI en haut de page",
    "i18n FR/EN - bascule locale localStorage",
    "Mode demo si backend inaccessible (3s)",
    "Inter + JetBrains Mono, prefers-reduced-motion",
], 7.15, 3.8, 5.7, size=12, spacing=0.44)
progress_bar(s8, 8, SEC["dash"])

# ═══════════════════════════════════════════════════════════════════
# Slide 9 — SECURITE  (Full statement + list)
# ═══════════════════════════════════════════════════════════════════
s9 = new_slide()
section_label(s9, "securite", SEC["secu"])
slide_num(s9, 9, SEC["secu"])
add_logo(s9)

headline(s9, "Production-ready\ndepuis le jour 1", t=0.5, size=40)

sec_items = [
    ("helmet",             "HTTP headers de securite (CSP, HSTS, X-Frame, nosniff)"),
    ("express-rate-limit", "100 req / 15min / IP - anti brute-force"),
    ("zod",                "Validation stricte des variables d'environnement au demarrage"),
    ("CORS strict",        "Whitelist explicite - pas de wildcard en production"),
    ("Non-root Docker",    "Tous les conteneurs en appuser (UID 1000)"),
    ("Secrets via env",    "Zero secret en dur - .env gitignored, .dockerignore complet"),
]
rect(s9, 0.4, 3.1, 12.6, 0.38, CARD)
txt(s9, "Composant", 0.6, 3.18, 3.5, 0.28, size=10, bold=True, color=SEC["secu"])
txt(s9, "Protection", 4.3, 3.18, 8.5, 0.28, size=10, bold=True, color=SEC["secu"])
for i, (comp, desc) in enumerate(sec_items):
    y = 3.52 + i*0.57
    bg = DARK_CARD if i%2==0 else CARD
    rect(s9, 0.4, y, 12.6, 0.54, bg)
    txt(s9, comp, 0.6,  y+0.1, 3.5, 0.36, size=12, bold=True, color=WHITE)
    txt(s9, desc, 4.3, y+0.1, 8.5, 0.36, size=11, color=MUTED)
progress_bar(s9, 9, SEC["secu"])

# ═══════════════════════════════════════════════════════════════════
# Slide 10 — i18n  (Split)
# ═══════════════════════════════════════════════════════════════════
s10 = new_slide()
section_label(s10, "internationalisation", SEC["i18n"])
slide_num(s10, 10, SEC["i18n"])
add_logo(s10)

headline(s10, "i18n FR/EN\nvanilla JS, zero dep", t=0.5, size=38)

rect(s10, 0.4, 3.2, 5.9, 3.9, DARK_CARD)
div10 = s10.shapes.add_shape(1, Inches(0.4), Inches(3.2), Inches(0.07), Inches(3.9))
div10.fill.solid(); div10.fill.fore_color.rgb = SEC["i18n"]; div10.line.fill.background()
txt(s10, "Architecture", 0.65, 3.3, 5.5, 0.4, size=13, bold=True, color=WHITE)
bullet(s10, [
    "Module IIFE I18n - zero dependance",
    "JSON async fetch depuis i18n/",
    "Persistence localStorage (cmm-locale)",
    "data-i18n sur tous les elements",
    "Intl API pour formatage locale-aware",
    "aria-pressed sur les boutons FR/EN",
], 0.65, 3.8, 5.5, size=12, spacing=0.44)

rect(s10, 6.9, 3.2, 6.1, 3.9, DARK_CARD)
div10b = s10.shapes.add_shape(1, Inches(6.9), Inches(3.2), Inches(0.07), Inches(3.9))
div10b.fill.solid(); div10b.fill.fore_color.rgb = GREEN; div10b.line.fill.background()
txt(s10, "Exemple fr.json", 7.15, 3.3, 5.7, 0.4, size=13, bold=True, color=WHITE)
txt(s10, '{\n  "title": "CRYPTO MONITOR",\n  "live": "EN DIRECT",\n  "connected": "Connecte",\n  "anomalies": {\n    "title": "Alertes anomalies",\n    "zscore": "Score Z"\n  },\n  "footer": {\n    "module": "Ingenierie Temps Reel"\n  }\n}',
    7.15, 3.85, 5.7, 3.0, size=10, color=GREEN)
progress_bar(s10, 10, SEC["i18n"])

# ═══════════════════════════════════════════════════════════════════
# Slide 11 — DEPLOIEMENT  (3 columns)
# ═══════════════════════════════════════════════════════════════════
s11 = new_slide()
section_label(s11, "deploiement", SEC["deploy"])
slide_num(s11, 11, SEC["deploy"])
add_logo(s11)

headline(s11, "3 options\nde deploiement", t=0.5, size=40)

cols = [
    ("01", "Lanceur\nautonomne",  ROSE,  "node launcher/src/index.mjs\n\nVerifie Docker, build, attend\nles healthchecks.\n\nCompile en .exe (pkg) :\nnpm run build:win"),
    ("02", "Image\ntout-en-un",   BLUE,  "docker build \\\n  -f all-in-one.Dockerfile \\\n  -t crypto-monitor .\n\ndocker run -p 8080:8080\n  crypto-monitor"),
    ("03", "Docker\nCompose",     PURPLE,"cp .env.example .env\ndocker-compose up -d\n\nDashboard  :8080\nKafka UI   :8090\nAPI        :3001"),
]
for i, (num, title, col, body_text) in enumerate(cols):
    x = 0.35 + i*4.35
    rect(s11, x, 3.15, 4.15, 4.0, DARK_CARD)
    top_bar = s11.shapes.add_shape(1, Inches(x), Inches(3.15), Inches(4.15), Inches(0.07))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = col; top_bar.line.fill.background()
    txt(s11, num,   x+0.18, 3.25, 3.8, 0.5, size=28, bold=True, color=col)
    txt(s11, title, x+0.18, 3.8,  3.8, 0.65, size=14, bold=True, color=WHITE)
    txt(s11, body_text, x+0.18, 4.55, 3.75, 2.5, size=10, color=MUTED)
progress_bar(s11, 11, SEC["deploy"])

# ═══════════════════════════════════════════════════════════════════
# Slide 12 — BILAN  (Full statement + points)
# ═══════════════════════════════════════════════════════════════════
s12 = new_slide()
section_label(s12, "bilan & perspectives", SEC["bilan"])
slide_num(s12, 12, SEC["bilan"])
add_logo(s12)

headline(s12, "Realise a 2\npour une consigne de 4", t=0.5, size=38)

rect(s12, 0.4, 3.2, 5.9, 3.9, DARK_CARD)
div12 = s12.shapes.add_shape(1, Inches(0.4), Inches(3.2), Inches(0.07), Inches(3.9))
div12.fill.solid(); div12.fill.fore_color.rgb = ROSE; div12.line.fill.background()
txt(s12, "Livre", 0.65, 3.3, 5.5, 0.4, size=13, bold=True, color=WHITE)
bullet(s12, [
    "Pipeline Kafka bout-en-bout fonctionnel",
    "3 services TypeScript production-ready",
    "Dashboard EFREI-branded FR/EN + mode demo",
    "Securite : helmet, rate-limit, zod, non-root",
    "3 options deploiement dont .exe standalone",
], 0.65, 3.8, 5.5, size=12, spacing=0.44)

rect(s12, 6.9, 3.2, 6.1, 3.9, DARK_CARD)
div12b = s12.shapes.add_shape(1, Inches(6.9), Inches(3.2), Inches(0.07), Inches(3.9))
div12b.fill.solid(); div12b.fill.fore_color.rgb = BLUE; div12b.line.fill.background()
txt(s12, "Evolutions possibles", 7.15, 3.3, 5.7, 0.4, size=13, bold=True, color=WHITE)
bullet(s12, [
    "SOL, BNB, XRP - nouvelles paires",
    "Alertes email/SMS via Kafka Connect",
    "InfluxDB / TimescaleDB pour historique",
    "Kubernetes + Helm pour la scalabilite",
    "Auth JWT REST + Socket.IO",
], 7.15, 3.8, 5.7, size=12, spacing=0.44)

progress_bar(s12, 12, ROSE)

prs.save(OUT)
print("PPTX v2 saved -> " + OUT)

